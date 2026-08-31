"""Generate an editable, boardroom-grade PowerPoint (.pptx) deck of a school's
external-exam performance (WAEC + JAMB together) — for administration, parent and
board meetings.

Data-driven from the same WAEC/JAMB analytics as the board-pack PDF and framed to
lead with the school's strengths. Every figure is aggregate (plus anonymous
distribution counts) — no student records. Charts are *native* PowerPoint charts
(not images), so the school can recolour, retitle and edit them like any deck they
built themselves in PowerPoint / Keynote / Google Slides."""
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK

# 16:9 canvas
_W = Inches(13.333)
_H = Inches(7.5)

_INK = RGBColor(0x14, 0x21, 0x1C)
_PRIMARY = RGBColor(0x0D, 0x6A, 0x4E)      # deep green
_PRIMARY2 = RGBColor(0x2E, 0x8B, 0x63)     # lighter green (2nd shade)
_ACCENT = RGBColor(0xC9, 0xA2, 0x27)       # gold
_NAVY = RGBColor(0x1F, 0x3A, 0x5F)         # JAMB accent
_RED = RGBColor(0xB4, 0x3A, 0x2E)          # weak / fail
_MUTED = RGBColor(0x6B, 0x7A, 0x74)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT = RGBColor(0xF4, 0xF7, 0xF5)
_GRID = RGBColor(0xE3, 0xE9, 0xE6)

_FONT = 'Calibri'


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _rect(slide, l, t, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _set(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=_FONT):
    p.text = str(text)
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return p


def _num(v, suffix=''):
    if v is None:
        return '—'
    try:
        f = float(v)
        return (f'{f:g}' if f != int(f) else f'{int(f)}') + suffix
    except (TypeError, ValueError):
        return str(v)


def _footer(slide, school, year):
    _rect(slide, 0, _H - Inches(0.32), _W, Inches(0.32), _LIGHT)
    tf = _box(slide, Inches(0.5), _H - Inches(0.34), Inches(12.3), Inches(0.3))
    _set(tf.paragraphs[0], f'{school}  ·  External Examination Performance {year}', 9, _MUTED)


# --------------------------------------------------------------------------- #
# Chart helpers (native pptx charts)
# --------------------------------------------------------------------------- #

def _style_axes(chart, number_format=None):
    try:
        cat = chart.category_axis
        cat.tick_labels.font.size = Pt(11)
        cat.tick_labels.font.name = _FONT
        cat.format.line.color.rgb = _GRID
        cat.major_tick_mark = XL_TICK_MARK.NONE
        cat.minor_tick_mark = XL_TICK_MARK.NONE
    except Exception:
        pass
    try:
        val = chart.value_axis
        val.tick_labels.font.size = Pt(10)
        val.tick_labels.font.name = _FONT
        val.has_major_gridlines = True
        val.major_gridlines.format.line.color.rgb = _GRID
        val.major_gridlines.format.line.width = Pt(0.5)
        val.format.line.fill.background()
        val.major_tick_mark = XL_TICK_MARK.NONE
        if number_format:
            val.tick_labels.number_format = number_format
            val.tick_labels.number_format_is_linked = False
    except Exception:
        pass


def _colour_points(series, colours):
    """Colour each bar/column individually (single-series charts read cleaner)."""
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colours[i % len(colours)]
        pt.format.line.fill.background()


def _add_bar(slide, l, t, w, h, categories, values, *, horizontal=False,
             point_colours=None, series_colour=_PRIMARY, number_format='0',
             data_labels=True, label_colour=_INK, label_size=10, max_scale=None):
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series('Series 1', list(values), number_format=number_format)
    ctype = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(ctype, l, t, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.vary_by_categories = False
    s = plot.series[0]
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = series_colour
    s.format.line.fill.background()
    if point_colours:
        _colour_points(s, point_colours)
    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = number_format
        dl.number_format_is_linked = False
        dl.font.size = Pt(label_size)
        dl.font.bold = True
        dl.font.name = _FONT
        dl.font.color.rgb = label_colour
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    _style_axes(chart, number_format)
    if max_scale is not None:
        try:
            chart.value_axis.maximum_scale = max_scale
            chart.value_axis.minimum_scale = 0
        except Exception:
            pass
    return chart


def _add_donut(slide, l, t, w, h, pairs, colours):
    """pairs = [(label, value), ...]."""
    cd = CategoryChartData()
    cd.categories = [p[0] for p in pairs]
    cd.add_series('Share', [p[1] for p in pairs], number_format='0"%"')
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = _FONT
    plot = chart.plots[0]
    try:
        plot.donut_hole_size = 62
    except Exception:
        pass
    s = plot.series[0]
    for i, pt in enumerate(s.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colours[i % len(colours)]
        pt.format.line.color.rgb = _WHITE
        pt.format.line.width = Pt(1.5)
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0"%"'
    dl.number_format_is_linked = False
    dl.font.size = Pt(11)
    dl.font.bold = True
    dl.font.name = _FONT
    dl.font.color.rgb = _WHITE
    return chart


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #

def _title_slide(prs, school, year, subtitle, generated):
    s = _blank(prs)
    _rect(s, 0, 0, _W, _H, _PRIMARY)
    _rect(s, Inches(0.9), Inches(4.85), Inches(3.4), Inches(0.10), _ACCENT)
    tf = _box(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
    _set(tf.paragraphs[0], 'EXTERNAL EXAMINATION PERFORMANCE', 15, _ACCENT, bold=True)
    tf2 = _box(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(2.4))
    _set(tf2.paragraphs[0], school, 44, _WHITE, bold=True)
    _set(tf2.add_paragraph(), f'WAEC & JAMB Results · {year}', 26, RGBColor(0xE6, 0xEF, 0xEA))
    if subtitle:
        _set(tf2.add_paragraph(), subtitle, 16, RGBColor(0xCF, 0xDD, 0xD7))
    ft = _box(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
    _set(ft.paragraphs[0], f'Board & Management Briefing  ·  Prepared {generated}',
         12, RGBColor(0xCF, 0xDD, 0xD7))


def _section_header(slide, kicker, title, subtitle=None):
    _rect(slide, 0, 0, _W, Inches(1.28), _PRIMARY)
    _rect(slide, 0, Inches(1.28), _W, Inches(0.06), _ACCENT)
    if kicker:
        kf = _box(slide, Inches(0.6), Inches(0.16), Inches(12), Inches(0.35))
        _set(kf.paragraphs[0], kicker.upper(), 11, _ACCENT, bold=True)
    tf = _box(slide, Inches(0.6), Inches(0.46), Inches(12.1), Inches(0.75))
    _set(tf.paragraphs[0], title, 26, _WHITE, bold=True)
    if subtitle:
        st = _box(slide, Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.5))
        _set(st.paragraphs[0], subtitle, 13, _MUTED)


def _stat_card(slide, l, t, w, value, label, accent=_ACCENT, sub=None):
    _rect(slide, l, t, w, Inches(1.95), _LIGHT)
    _rect(slide, l, t, w, Inches(0.10), accent)
    vf = _box(slide, l, t + Inches(0.30), w, Inches(0.95))
    vf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(vf.paragraphs[0], value, 38, _PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    lf = _box(slide, l, t + Inches(1.28), w, Inches(0.6))
    _set(lf.paragraphs[0], label, 11.5, _MUTED, align=PP_ALIGN.CENTER)
    if sub:
        _set(lf.add_paragraph(), sub, 9.5, accent, bold=True, align=PP_ALIGN.CENTER)


def _kpi_slide(prs, school, year, waec, jamb, cutoff, correlation):
    s = _blank(prs)
    _section_header(s, 'Executive summary', 'Performance at a glance',
                    'Headline WAEC & JAMB outcomes for the year')
    cards = []
    if waec:
        cards.append((_num(waec.get('overall_pass_rate'), '%'), 'WAEC credit/pass rate', _PRIMARY, 'WAEC'))
        cards.append((_num(waec.get('overall_distinction_rate'), '%'), 'WAEC A1–B3 distinctions', _ACCENT, 'WAEC'))
    if jamb:
        cards.append((_num(jamb.get('mean_score')), 'JAMB mean score (of 400)', _NAVY, 'JAMB'))
    if cutoff:
        cards.append((_num(cutoff.get('eligible_200_pct'), '%'), 'JAMB candidates ≥ 200', _NAVY, 'JAMB'))
    if len(cards) < 4 and jamb:
        cards.append((_num(jamb.get('max_score')), 'JAMB highest score', _NAVY, 'JAMB'))
    cards = cards[:4] or [('—', 'No data', _MUTED, None)]
    n = len(cards)
    gap = Inches(0.4)
    total_w = _W - Inches(1.2) - gap * (n - 1)
    w = Emu(int(total_w / n))
    x = Inches(0.6)
    for value, label, accent, sub in cards:
        _stat_card(s, x, Inches(2.3), w, value, label, accent, sub)
        x = Emu(int(x) + int(w) + int(gap))
    parts = []
    if waec and jamb:
        parts.append(f"{_num(waec.get('unique_students'))} WAEC candidates and "
                     f"{_num(jamb.get('total_students'))} JAMB candidates assessed.")
    if correlation and not correlation.get('error') and correlation.get('predictive_power'):
        parts.append(f"WAEC↔JAMB link: {correlation.get('predictive_power')}.")
    if parts:
        nf = _box(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.6))
        _set(nf.paragraphs[0], '  '.join(parts), 14, _INK)
    _footer(s, school, year)


def _waec_grade_slide(prs, school, year, waec):
    dist = (waec or {}).get('grade_distribution') or {}
    order = ['A1', 'B2', 'B3', 'C4', 'C5', 'C6', 'D7', 'E8', 'F9']
    cats = [g for g in order if dist.get(g)]
    if not cats:
        return
    vals = [dist[g] for g in cats]

    def col(g):
        if g in ('E8', 'F9'):
            return _RED
        if g == 'D7':
            return _ACCENT
        return _PRIMARY if g in ('A1', 'B2', 'B3') else _PRIMARY2

    s = _blank(prs)
    _section_header(s, 'WAEC', 'Grade distribution',
                    'How every subject entry was graded — the shape of the results')
    _add_bar(s, Inches(0.7), Inches(1.7), Inches(8.4), Inches(5.2), cats, vals,
             point_colours=[col(g) for g in cats], number_format='0')
    total = sum(vals)
    credits = sum(v for g, v in zip(cats, vals) if g in ('A1', 'B2', 'B3', 'C4', 'C5', 'C6'))
    distinctions = sum(v for g, v in zip(cats, vals) if g in ('A1', 'B2', 'B3'))
    lf = _box(s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(4.5))
    _set(lf.paragraphs[0], 'READING THE CHART', 11, _ACCENT, bold=True)
    for label, val in [('Total subject entries', total),
                       ('Credits (A1–C6)', f'{credits}  ({round(credits/total*100)}%)' if total else '—'),
                       ('Distinctions (A1–B3)', f'{distinctions}  ({round(distinctions/total*100)}%)' if total else '—')]:
        p = lf.add_paragraph(); _set(p, label, 12, _MUTED); p.space_before = Pt(10)
        _set(lf.add_paragraph(), str(val), 18, _PRIMARY, bold=True)
    _footer(s, school, year)


def _waec_subjects_slide(prs, school, year, waec):
    subs = (waec or {}).get('subject_analysis') or []
    if not subs:
        return
    top = sorted(subs, key=lambda x: x.get('pass_rate') or 0, reverse=True)[:10]
    cats = [x['subject'] for x in top][::-1]           # reversed: strongest on top of a bar chart
    vals = [x.get('pass_rate') or 0 for x in top][::-1]
    cols = [(_PRIMARY if v >= 75 else _ACCENT if v >= 50 else _RED) for v in vals]
    s = _blank(prs)
    _section_header(s, 'WAEC', 'Subject pass rates',
                    'Share of entries at credit or better, strongest subjects first')
    _add_bar(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.2), cats, vals,
             horizontal=True, point_colours=cols, number_format='0"%"', max_scale=100)
    _footer(s, school, year)


def _jamb_distribution_slide(prs, school, year, jamb):
    dist = (jamb or {}).get('distribution') or {}
    cats = [k for k in ['0-100', '101-150', '151-200', '201-250', '251-300',
                        '301-350', '351-400'] if k in dist]
    if not cats or not sum(dist.get(k, 0) for k in cats):
        return
    vals = [dist[k] for k in cats]

    def col(band):
        low = int(band.split('-')[0])
        return _PRIMARY if low >= 250 else _PRIMARY2 if low >= 200 else _NAVY if low >= 150 else _RED

    s = _blank(prs)
    _section_header(s, 'JAMB', 'Score distribution',
                    'Candidates by score band (of 400) — the university-admission picture')
    _add_bar(s, Inches(0.7), Inches(1.7), Inches(8.4), Inches(5.2), cats, vals,
             point_colours=[col(b) for b in cats], number_format='0')
    lf = _box(s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(4.6))
    _set(lf.paragraphs[0], 'KEY THRESHOLDS', 11, _ACCENT, bold=True)
    for label, val in [('Mean score', _num(jamb.get('mean_score'))),
                       ('Median score', _num(jamb.get('median_score'))),
                       ('Highest score', _num(jamb.get('max_score'))),
                       ('Std. deviation', _num(jamb.get('std_deviation')))]:
        p = lf.add_paragraph(); _set(p, label, 12, _MUTED); p.space_before = Pt(8)
        _set(lf.add_paragraph(), str(val), 18, _NAVY, bold=True)
    _footer(s, school, year)


def _readiness_slide(prs, school, year, jamb, cutoff):
    if not cutoff:
        return
    bands = [('≥ 200 Admissible', cutoff.get('eligible_200'), cutoff.get('eligible_200_pct'), _PRIMARY2),
             ('≥ 250 Competitive', cutoff.get('competitive_250'), cutoff.get('competitive_250_pct'), _PRIMARY),
             ('≥ 300 Elite', cutoff.get('elite_300'), cutoff.get('elite_300_pct'), _ACCENT)]
    if not any(b[1] for b in bands):
        return
    s = _blank(prs)
    _section_header(s, 'JAMB', 'University readiness',
                    'How many candidates clear each admission threshold')
    cats = [b[0] for b in bands]
    vals = [b[1] or 0 for b in bands]
    _add_bar(s, Inches(0.7), Inches(1.8), Inches(7.3), Inches(4.9), cats, vals,
             point_colours=[b[3] for b in bands], number_format='0')
    y = Inches(2.0)
    for label, cnt, pct, colr in bands:
        _rect(s, Inches(8.4), y, Inches(4.2), Inches(1.35), _LIGHT)
        _rect(s, Inches(8.4), y, Inches(0.10), Inches(1.35), colr)
        tf = _box(s, Inches(8.7), y + Inches(0.12), Inches(3.7), Inches(1.15))
        _set(tf.paragraphs[0], f'{_num(pct, "%")} of candidates', 20, colr, bold=True)
        _set(tf.add_paragraph(), f'{label}  ·  {_num(cnt)} students', 12, _MUTED)
        y = Emu(int(y) + int(Inches(1.55)))
    _footer(s, school, year)


def _combined_slide(prs, school, year, waec, jamb, cutoff):
    """One slide showing WAEC and JAMB side by side — the 'together' view."""
    if not (waec and jamb):
        return
    s = _blank(prs)
    _section_header(s, 'WAEC + JAMB', 'The two exams together',
                    'Certificate outcomes and university-entrance outcomes, side by side')
    _set(_box(s, Inches(0.9), Inches(1.7), Inches(5), Inches(0.4)).paragraphs[0],
         'WAEC — credit outcomes', 14, _PRIMARY, bold=True)
    pr = waec.get('overall_pass_rate') or 0
    _add_donut(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.4),
               [('Credit / pass', round(pr, 1)), ('Below credit', round(100 - pr, 1))],
               [_PRIMARY, _GRID])
    _set(_box(s, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.4)).paragraphs[0],
         'JAMB — university admissible (≥ 200)', 14, _NAVY, bold=True)
    e200 = (cutoff or {}).get('eligible_200_pct') or 0
    _add_donut(s, Inches(7.1), Inches(2.1), Inches(5.6), Inches(4.4),
               [('Admissible ≥ 200', round(e200, 1)), ('Below 200', round(100 - e200, 1))],
               [_NAVY, _GRID])
    _footer(s, school, year)


def _table_slide(prs, school, year, kicker, title, subtitle, headers, rows):
    s = _blank(prs)
    _section_header(s, kicker, title, subtitle)
    if not rows:
        _set(_box(s, Inches(0.6), Inches(2.5), Inches(11), Inches(1)).paragraphs[0],
             'No data available for this section.', 16, _MUTED)
        _footer(s, school, year)
        return s
    rows = rows[:10]
    nrows, ncols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.7),
                            _W - Inches(1.2), Inches(0.46) * nrows).table
    for c, h in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = str(h)
        cell.fill.solid(); cell.fill.fore_color.rgb = _PRIMARY
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.color.rgb = _WHITE
        p.runs[0].font.bold = True; p.runs[0].font.size = Pt(13); p.runs[0].font.name = _FONT
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = _INK
            p.runs[0].font.name = _FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = _LIGHT if r % 2 == 0 else _WHITE
    _footer(s, school, year)
    return s


def _insights_slide(prs, school, year, insights):
    s = _blank(prs)
    _section_header(s, 'Analysis', 'Key takeaways', 'What the results tell us')
    colours = {'positive': _PRIMARY, 'warning': _ACCENT, 'negative': _RED, 'critical': _RED}
    tf = _box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5.0))
    first = True
    for i in (insights or [])[:6]:
        colr = colours.get(i.get('level'), _PRIMARY)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(10)
        first = False
        _set(p, '▪  ' + i.get('title', ''), 17, colr, bold=True)
        detail = (i.get('detail') or '')
        if detail:
            dp = tf.add_paragraph()
            _set(dp, '     ' + detail, 12.5, _MUTED)
    if first:
        _set(tf.paragraphs[0], 'Results are being compiled.', 16, _MUTED)
    _footer(s, school, year)


def _recommendations_slide(prs, school, year, waec, jamb, cutoff, correlation):
    """A business-style 'so what / now what' slide derived from the figures."""
    recs = []
    if waec and waec.get('subject_analysis'):
        weak = [w for w in sorted(waec['subject_analysis'], key=lambda x: x.get('pass_rate') or 0)[:3]
                if (w.get('pass_rate') or 0) < 75]
        if weak:
            names = ', '.join(w['subject'] for w in weak)
            recs.append(('Target intervention in weaker WAEC subjects',
                         f'Prioritise {names} — each below a 75% credit rate. Assign lead '
                         f'teachers, add clinic sessions and re-test at mid-term.'))
    if cutoff is not None:
        below = 100 - (cutoff.get('eligible_200_pct') or 0)
        if below > 20:
            recs.append(('Lift the JAMB admissible share above 200',
                         f'{round(below)}% of candidates scored below 200. Run focused CBT '
                         f'drills and past-question practice for the borderline cohort.'))
        if (cutoff.get('elite_300_pct') or 0) < 10:
            recs.append(('Grow the elite (≥ 300) band',
                         'Stretch top performers with timed mock JAMB and subject masterclasses '
                         'to convert strong candidates into 300+ scorers.'))
    if correlation and not correlation.get('error') and correlation.get('predictive_power'):
        recs.append(('Use WAEC as an early-warning signal',
                     f'WAEC and JAMB move together ({correlation.get("predictive_power")}). '
                     f'Flag students with weak WAEC mocks for JAMB support early.'))
    recs.append(('Share this pack with stakeholders',
                 'Brief staff, parents and the board with these figures, then agree owners and '
                 'a review date for each action above.'))
    s = _blank(prs)
    _section_header(s, 'Recommendations', 'From results to action',
                    'Priorities the data points to for the year ahead')
    tf = _box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(5.1))
    first = True
    for n, (title, detail) in enumerate(recs[:5], start=1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(12)
        first = False
        _set(p, f'{n}.  {title}', 16, _PRIMARY, bold=True)
        _set(tf.add_paragraph(), '     ' + detail, 12.5, _MUTED)
    _footer(s, school, year)


def _closing_slide(prs, school):
    s = _blank(prs)
    _rect(s, 0, 0, _W, _H, _PRIMARY)
    _rect(s, Inches(4.97), Inches(4.35), Inches(3.4), Inches(0.08), _ACCENT)
    tf = _box(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.5))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], 'Thank you', 42, _WHITE, bold=True, align=PP_ALIGN.CENTER)
    _set(tf.add_paragraph(), school, 20, RGBColor(0xE6, 0xEF, 0xEA), align=PP_ALIGN.CENTER)
    qf = _box(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(1.2))
    _set(qf.paragraphs[0], 'Questions & discussion', 15, RGBColor(0xCF, 0xDD, 0xD7),
         align=PP_ALIGN.CENTER)


def build_deck(*, year, school_name, generated, branch_label=None,
               waec_stats=None, jamb_stats=None, cutoff=None, correlation=None, insights=None):
    """Return the .pptx bytes for the external-exam (WAEC + JAMB) results deck."""
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    school = school_name or 'Our School'
    subtitle = branch_label or 'All branches'

    _title_slide(prs, school, year, subtitle, generated)
    _kpi_slide(prs, school, year, waec_stats, jamb_stats, cutoff, correlation)

    # WAEC section
    _waec_grade_slide(prs, school, year, waec_stats)
    _waec_subjects_slide(prs, school, year, waec_stats)

    # JAMB section
    _jamb_distribution_slide(prs, school, year, jamb_stats)
    _readiness_slide(prs, school, year, jamb_stats, cutoff)

    # Combined view + detail tables
    _combined_slide(prs, school, year, waec_stats, jamb_stats, cutoff)

    if jamb_stats and jamb_stats.get('subject_analysis'):
        jrows = [(x['subject'], _num(x.get('mean_score')), _num(x.get('above_50')),
                  _num(x.get('above_70')))
                 for x in jamb_stats['subject_analysis'][:10]]
        _table_slide(prs, school, year, 'JAMB', 'Subject performance',
                     'Mean score and candidates clearing 50 / 70 per subject',
                     ['Subject', 'Mean', '≥ 50', '≥ 70'], jrows)

    if jamb_stats:
        jkpi = [('Candidates', _num(jamb_stats.get('total_students'))),
                ('Mean score', _num(jamb_stats.get('mean_score'))),
                ('Median score', _num(jamb_stats.get('median_score'))),
                ('Highest score', _num(jamb_stats.get('max_score'))),
                ('Lowest score', _num(jamb_stats.get('min_score'))),
                ('Scored ≥ 200', _num(jamb_stats.get('above_200'))),
                ('Scored ≥ 250', _num(jamb_stats.get('above_250'))),
                ('Scored ≥ 300', _num(jamb_stats.get('above_300')))]
        _table_slide(prs, school, year, 'JAMB', 'Cohort statistics',
                     'The full JAMB picture for the year', ['Metric', 'Value'], jkpi)

    _insights_slide(prs, school, year, insights)
    _recommendations_slide(prs, school, year, waec_stats, jamb_stats, cutoff, correlation)
    _closing_slide(prs, school)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
