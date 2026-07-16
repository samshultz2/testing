"""External-exam PowerPoint deck: the generator builds a valid, data-driven
.pptx, and the download route serves it (gated to a year with data)."""
import io

from config import Config
from tests.conftest import login_token


def _admin(app):
    c = app.test_client()
    c.post('/login', data={'password': Config.ADMIN_PASSWORD, '_csrf_token': login_token(c)})
    return c


def _sample():
    waec = {'overall_pass_rate': 82.5, 'overall_distinction_rate': 41.0, 'unique_students': 120,
            'subject_analysis': [{'subject': 'Mathematics', 'pass_rate': 88},
                                 {'subject': 'Biology', 'pass_rate': 91}]}
    jamb = {'total_students': 95, 'mean_score': 241, 'max_score': 312,
            'above_200': 70, 'above_250': 38, 'above_300': 6}
    cutoff = {'eligible_200_pct': 73.7, 'competitive_250_pct': 40.0, 'elite_300_pct': 6.3}
    insights = [{'level': 'positive', 'title': 'Strong WAEC outcomes', 'detail': '82.5% pass rate.'}]
    return waec, jamb, cutoff, insights


def test_build_deck_is_valid_pptx_with_data():
    from pptx import Presentation
    from utils.exam_deck import build_deck
    waec, jamb, cutoff, insights = _sample()
    data = build_deck(year=2025, school_name='Greenfield Academy', generated='2025-07-16',
                      branch_label='Main', waec_stats=waec, jamb_stats=jamb,
                      cutoff=cutoff, insights=insights)
    assert data[:2] == b'PK'                          # a real Office Open XML (zip) file
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    assert len(slides) >= 5
    blob = []
    for sl in slides:
        for sh in sl.shapes:
            if sh.has_text_frame:
                blob.append(sh.text_frame.text)
            if sh.has_table:
                blob += [c.text for row in sh.table.rows for c in row.cells]
    text = ' | '.join(blob)
    assert 'Greenfield Academy' in text               # school named
    assert '82.5%' in text and '241' in text          # WAEC + JAMB figures
    assert 'Strong WAEC outcomes' in text             # insights carried through


def test_deck_survives_missing_sections():
    from pptx import Presentation
    from utils.exam_deck import build_deck
    # only JAMB data present — must still produce a valid deck, no crash
    data = build_deck(year=2024, school_name='X School', generated='2024-01-01',
                      jamb_stats={'total_students': 10, 'mean_score': 190})
    assert Presentation(io.BytesIO(data)) is not None


def test_deck_route_requires_year(app):
    c = _admin(app)
    r = c.get('/results/analytics/deck.pptx', follow_redirects=False)
    assert r.status_code in (302, 303)                # no year -> redirect back


def test_deck_route_serves_pptx_when_data_exists(app):
    from models import db, WAECResult, Student, Branch
    with app.app_context():
        from utils.finance_ledger import ensure_tables; ensure_tables()
        st = Student(student_id=Student.generate_student_id(), first_name='A', surname='B',
                     gender='Male', is_active=True, branch_id=Branch.get_default().id)
        db.session.add(st); db.session.commit()
        db.session.add(WAECResult(student_id=st.id, exam_year=2025, subject='Mathematics', grade='B2'))
        db.session.commit()
    c = _admin(app)
    r = c.get('/results/analytics/deck.pptx?year=2025')
    assert r.status_code == 200
    assert 'presentationml' in r.headers.get('Content-Type', '')
    assert '.pptx' in r.headers.get('Content-Disposition', '')
    assert r.get_data()[:2] == b'PK'
